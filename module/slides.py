
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_slide(prs, title, content):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    textbox = slide.placeholders[1]
    textbox.text = content
    for paragraph in textbox.text_frame.paragraphs:
        paragraph.space_after = Pt(10)
        paragraph.alignment = PP_ALIGN.LEFT

def add_image_slide(prs, title, image_path):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    left = Inches(0.75)
    top = Inches(1.25)
    slide.shapes.add_picture(image_path, left, top, width=Inches(8.5))

def generate_presentation(content_dict, output_file="LangChain_Project_Slides.pptx", diagram_path=None):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    
    # Explicitly set title
    slide.shapes.title.text = content_dict["title"]

    # Clear subtitle if present
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = ""

    # Add the rest of the slides
    for section, text in content_dict.items():
        if section == "title":
            continue
        create_slide(prs, section, text)

    if diagram_path:
        add_image_slide(prs, "Solution Architecture Diagram", diagram_path)

    prs.save(output_file)
    print(f"Presentation saved as {output_file}")


    if diagram_path:
        add_image_slide(prs, "Solution Architecture Diagram", diagram_path)

    prs.save(output_file)
    print(f"Presentation saved as {output_file}")
